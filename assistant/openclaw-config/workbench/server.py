#!/usr/bin/env python3
"""
Bolt Console（OpenClaw）本地后端
- 仅绑定 127.0.0.1，个人本机使用
- 前端通过 fetch 调用以下 API
- 核心路径无第三方依赖；Office 预览优先用 ~/.openclaw/venv
"""
import html as htmlmod
import json
import os
import queue
import re
import stat
import subprocess
import threading
import time
import urllib.parse
import zipfile
import xml.etree.ElementTree as ET
from http.server import BaseHTTPRequestHandler, ThreadingHTTPServer
from pathlib import Path

PORT = 18790
HOST = "127.0.0.1"
PRODUCT = "Bolt Console"
VERSION = "2.0"

HOME = str(Path.home())
TASKS_DIR = os.path.join(HOME, "tasks")
WORKSPACE_DIR = os.path.join(HOME, ".openclaw", "workspace")
MEMORY_DIR = os.path.join(WORKSPACE_DIR, "memory")
SKILLS_DIR = os.path.join(WORKSPACE_DIR, "skills")
WORKBENCH_DIR = os.path.dirname(os.path.abspath(__file__))
HISTORY_FILE = os.path.join(WORKBENCH_DIR, "history.json")
SESSIONS_FILE = os.path.join(WORKBENCH_DIR, "sessions.json")
VENV_PY = os.path.join(HOME, ".openclaw", "venv", "bin", "python3")
PANDOC = os.path.join(
    HOME, ".openclaw", "venv", "lib", "python3.9",
    "site-packages", "pypandoc", "files", "pandoc",
)
CURSOR_RUNNER = os.path.join(WORKBENCH_DIR, "cursor", "cursor-agent.mjs")
CURSOR_CWD = os.environ.get("BOLT_CURSOR_CWD") or TASKS_DIR

ALLOWED_ROOTS = [TASKS_DIR, MEMORY_DIR, SKILLS_DIR]

# region agent log
DEBUG_LOG = "/Users/chingching/ai-assistant/.cursor/debug-019ce7.log"

def _dbg(message, hypothesis_id, **data):
    try:
        rec = {
            "sessionId": "019ce7",
            "location": "server.py",
            "message": message,
            "hypothesisId": hypothesis_id,
            "data": data,
            "timestamp": int(time.time() * 1000),
        }
        with open(DEBUG_LOG, "a", encoding="utf-8") as f:
            f.write(json.dumps(rec, ensure_ascii=False) + "\n")
    except Exception:
        pass
# endregion

IMAGE_EXTS = {".png", ".jpg", ".jpeg", ".gif", ".webp", ".svg", ".bmp", ".ico"}
AUDIO_EXTS = {".mp3", ".wav", ".m4a", ".aac", ".ogg", ".flac", ".opus"}
VIDEO_EXTS = {".mp4", ".webm", ".mov", ".m4v"}
MD_EXTS = {".md", ".markdown"}
PDF_EXTS = {".pdf"}
OFFICE_PREVIEW_EXTS = {".docx", ".pptx", ".xlsx"}
OFFICE_DOWNLOAD_EXTS = {".doc", ".xls", ".ppt", ".odt", ".ods", ".odp"}
TEXT_EXTS = {
    ".md", ".markdown", ".txt", ".html", ".json", ".csv", ".log", ".py", ".sh",
    ".mjs", ".js", ".css", ".yml", ".yaml", ".xml", ".ts", ".toml", ".ini", ".lrc",
}

MIME = {
    ".png": "image/png", ".jpg": "image/jpeg", ".jpeg": "image/jpeg",
    ".gif": "image/gif", ".webp": "image/webp", ".svg": "image/svg+xml",
    ".bmp": "image/bmp", ".ico": "image/x-icon",
    ".mp3": "audio/mpeg", ".wav": "audio/wav", ".m4a": "audio/mp4",
    ".aac": "audio/aac", ".ogg": "audio/ogg", ".flac": "audio/flac",
    ".opus": "audio/opus",
    ".mp4": "video/mp4", ".webm": "video/webm", ".mov": "video/quicktime",
    ".m4v": "video/mp4",
    ".pdf": "application/pdf",
}

AGENTS = ["main", "coder", "office", "cursor"]

_tasks = {}
_hist = []
_sessions = []  # 会话列表，元素：{session_id, agent, agent_id, title, created, updated, messages:[]}
_lock = threading.Lock()
# task_id -> list[queue.Queue]：SSE 订阅者，cursor 任务运行时会向其中推送流式事件
_stream_queues = {}
# task_id -> subprocess.Popen：运行中的子进程，用于「停止」
_procs = {}
# task_id -> threading.Event：取消标记
_cancel = {}


def _publish(task_id, ev):
    """把一条事件推给该任务的所有 SSE 订阅者（无订阅者则忽略）。"""
    subs = _stream_queues.get(task_id)
    if not subs:
        return
    for q in list(subs):
        try:
            q.put(ev)
        except Exception:  # noqa: BLE001
            pass


def log(msg):
    print("[console] %s" % msg, flush=True)


def _is_cancelled(task_id):
    ev = _cancel.get(task_id)
    return bool(ev and ev.is_set())


def cancel_task(task_id):
    """停止一个 running 任务：杀死子进程、标记 cancelled、推送事件。返回是否真的取消。"""
    with _lock:
        t = _tasks.get(task_id)
        if not t or t.get("status") != "running":
            return False
        _cancel.setdefault(task_id, threading.Event()).set()
        proc = _procs.get(task_id)
        if proc is not None and proc.poll() is None:
            try:
                proc.kill()
            except Exception:  # noqa: BLE001
                pass
        t["status"] = "cancelled"
        t["error"] = None
        t["activity"] = "已停止"
        t["duration"] = round(time.time() - t.get("started", time.time()), 1)
        session_id = t.get("session_id")
        _set_hist(task_id, "cancelled", None)
        if session_id:
            _append_session_msg(session_id, "assistant", "（用户已停止）", "cancelled")
    _publish(task_id, {"type": "cancelled"})
    _dbg("task cancelled", "D", task_id=task_id)
    return True


def file_kind(name):
    ext = os.path.splitext(name)[1].lower()
    if ext in IMAGE_EXTS:
        return "image"
    if ext in AUDIO_EXTS:
        return "audio"
    if ext in VIDEO_EXTS:
        return "video"
    if ext in MD_EXTS:
        return "markdown"
    if ext in PDF_EXTS:
        return "pdf"
    if ext in OFFICE_PREVIEW_EXTS or ext in OFFICE_DOWNLOAD_EXTS:
        return "office"
    if ext in TEXT_EXTS:
        return "text"
    return "other"


def mime_for(name):
    ext = os.path.splitext(name)[1].lower()
    return MIME.get(ext, "application/octet-stream")


# ---------------- 任务执行 ----------------

def run_openclaw_agent(task_id, agent, message, session_id=None):
    env = dict(os.environ)
    cmd = ["openclaw", "agent", "--agent", agent, "--json", "-m", message]
    started = time.time()
    log("run task %s agent=%s msg=%s..." % (task_id, agent, message[:60]))

    def _act(msg):
        """更新任务 activity 并推 SSE，让前端不再停留在「排队中…」。"""
        with _lock:
            t = _tasks.get(task_id)
            if t:
                t["activity"] = msg
        _publish(task_id, {"type": "status", "message": msg})

    try:
        proc = subprocess.Popen(
            cmd, stdout=subprocess.PIPE, stderr=subprocess.PIPE, text=True, env=env)
        with _lock:
            _procs[task_id] = proc
        _act("正在执行…")
        try:
            out, err = proc.communicate(timeout=900)
        except subprocess.TimeoutExpired:
            proc.kill()
            out, err = proc.communicate()
        finally:
            with _lock:
                _procs.pop(task_id, None)
        # 被手动停止：cancel_task 已标记 cancelled，这里不再覆盖
        if _is_cancelled(task_id):
            return
        text = None
        try:
            data = json.loads(out)
            payloads = (data.get("result") or {}).get("payloads") or []
            joined = "".join(p.get("text") or "" for p in payloads).strip()
            if joined:
                text = joined
            else:
                text = data.get("status") or data.get("summary") or out[-2000:]
        except json.JSONDecodeError:
            text = out.strip() or ("stderr: " + (err or "")[-1000:]) or "（无输出）"
        with _lock:
            _tasks[task_id] = {
                "task_id": task_id,
                "status": "done",
                "text": text,
                "error": None,
                "agent": agent,
                "message": message,
                "started": started,
                "duration": round(time.time() - started, 1),
            }
            _set_hist(task_id, "done", text)
            if session_id:
                _append_session_msg(session_id, "assistant", text, "done")
        _publish(task_id, {"type": "done", "text": text})
        # region agent log
        _dbg("task done", "B", task_id=task_id, text_len=len(text or ""))
        # endregion
    except Exception as e:  # noqa: BLE001
        if _is_cancelled(task_id):
            return
        log("task %s error: %s" % (task_id, e))
        with _lock:
            _tasks[task_id] = {
                "task_id": task_id,
                "status": "error",
                "text": "",
                "error": str(e),
                "agent": agent,
                "message": message,
                "started": started,
                "duration": round(time.time() - started, 1),
            }
            _set_hist(task_id, "error", None)
            if session_id:
                _append_session_msg(session_id, "assistant", str(e), "error")
        _publish(task_id, {"type": "error", "error": str(e)})
        # region agent log
        _dbg("task error", "B", task_id=task_id, error=str(e))
        # endregion


def run_cursor_agent(task_id, agent, message, session_id=None):
    """通过 Node 侧车调用 Cursor SDK，并实时解析其 NDJSON 流式输出。"""
    env = dict(os.environ)
    started = time.time()
    text_parts = []      # 累积的增量文本
    activity = "启动中…"
    log_lines = []       # 最近的活动日志（status/tool）
    final = None

    def _apply(final_text=None, final_status=None, err=None):
        with _lock:
            t = _tasks.get(task_id)
            if not t:
                return
            t["activity"] = activity
            t["log"] = list(log_lines[-40:])
            if final_status is not None:
                t["status"] = final_status
                t["text"] = final_text or "".join(text_parts)
                t["error"] = err
                t["duration"] = round(time.time() - started, 1)
            else:
                t["text"] = "".join(text_parts)

    def _finalize(ev):
        """把终态（done/error）一次性落盘：任务状态、历史、会话消息、agentId。"""
        et = ev.get("type")
        is_done = et == "done"
        text = (ev.get("text") or "").strip() or "".join(text_parts)
        err = ev.get("error")
        status = "done" if is_done else "error"
        with _lock:
            t = _tasks.get(task_id)
            if t:
                t["status"] = status
                t["text"] = text if is_done else (t.get("text") or "")
                t["error"] = err
                t["activity"] = activity
                t["log"] = list(log_lines[-40:])
                t["duration"] = round(time.time() - started, 1)
            _set_hist(task_id, status, text if is_done else None)
            if session_id:
                _append_session_msg(session_id, "assistant", text if is_done else (err or "执行失败"), status)
                if is_done:
                    _set_session_agent_id(session_id, ev.get("agentId"))

    try:
        log("run cursor task %s msg=%s..." % (task_id, message[:60]))
        # 续接会话：把 cursor 的 agentId 传给侧车，侧车用 Agent.resume 恢复上下文
        resume_id = None
        if session_id:
            s = _session_by_id(session_id)
            resume_id = (s or {}).get("agent_id")
        payload = json.dumps({
            "message": message,
            "cwd": CURSOR_CWD,
            "sessionId": resume_id or "",
        })
        proc = subprocess.Popen(
            ["node", CURSOR_RUNNER],
            stdin=subprocess.PIPE, stdout=subprocess.PIPE,
            stderr=subprocess.PIPE, text=True, env=env)
        with _lock:
            _procs[task_id] = proc
        try:
            proc.stdin.write(payload)
        except BrokenPipeError:
            pass
        try:
            proc.stdin.close()
        except Exception:  # noqa: BLE001
            pass

        def _reader():
            nonlocal activity, final
            for line in proc.stdout:
                line = line.strip()
                if not line:
                    continue
                try:
                    ev = json.loads(line)
                except json.JSONDecodeError:
                    continue
                et = ev.get("type")
                if et == "status":
                    activity = ev.get("message") or activity
                    log_lines.append(activity)
                elif et == "agent":
                    # 尽早持久化 agentId，供下一轮续会话（避免 done 之后才回填的时序竞争）
                    aid = ev.get("agentId")
                    if session_id and aid:
                        with _lock:
                            _set_session_agent_id(session_id, aid)
                elif et == "thinking":
                    text = (ev.get("text") or "").strip()
                    activity = "思考中…"
                    # 思考文本按行拆入日志，便于右上角 console 逐行展示
                    for ln in text.splitlines():
                        ln = ln.strip()
                        if ln:
                            log_lines.append("· " + (ln if len(ln) <= 180 else ln[:180] + "…"))
                elif et == "tool":
                    name = ev.get("tool") or ""
                    if ev.get("phase") == "running":
                        activity = "执行: %s" % name
                        log_lines.append(activity)
                    else:
                        activity = "完成: %s" % name
                        out = str(ev.get("output") or "").strip()
                        log_lines.append("✓ %s %s" % (name, out[:160]))
                elif et == "delta":
                    text_parts.append(ev.get("text") or "")
                elif et in ("done", "error"):
                    final = ev
                    _finalize(ev)
                    _publish(task_id, ev)
                    continue
                _apply()
                _publish(task_id, ev)

        reader = threading.Thread(target=_reader, daemon=True)
        reader.start()
        try:
            proc.wait(timeout=900)
        except subprocess.TimeoutExpired:
            proc.kill()
            if _is_cancelled(task_id):
                with _lock:
                    _procs.pop(task_id, None)
                return
            _apply(final_status="error", err="Cursor 执行超时（900s）")
            with _lock:
                _set_hist(task_id, "error", None)
                if session_id:
                    _append_session_msg(session_id, "assistant", "Cursor 执行超时（900s）", "error")
            return
        reader.join(timeout=5)
        with _lock:
            _procs.pop(task_id, None)
        # 被手动停止：cancel_task 已标记 cancelled，直接返回，不覆盖
        if _is_cancelled(task_id):
            return
        if final is None:
            err = (proc.stderr.read() or "")[-1000:]
            _finalize({"type": "error", "error": "Cursor 无有效输出" + ("\n" + err if err else "")})
            _dbg("cursor task error", "C", task_id=task_id, error="no output")
            return
        if final.get("type") == "done":
            _dbg("cursor task done", "C", task_id=task_id,
                 text_len=len(final.get("text") or ""), agent_id=final.get("agentId"))
        else:
            _dbg("cursor task error", "C", task_id=task_id, error=final.get("error"))
    except Exception as e:  # noqa: BLE001
        if _is_cancelled(task_id):
            with _lock:
                _procs.pop(task_id, None)
            return
        log("cursor task %s error: %s" % (task_id, e))
        _apply(final_status="error", err=str(e))
        with _lock:
            _set_hist(task_id, "error", None)
            if session_id:
                _append_session_msg(session_id, "assistant", str(e), "error")
        # region agent log
        _dbg("cursor task error", "C", task_id=task_id, error=str(e))
        # endregion


def save_history():
    try:
        with open(HISTORY_FILE, "w", encoding="utf-8") as f:
            json.dump(_hist[-50:], f, ensure_ascii=False, indent=2)
    except Exception:  # noqa: BLE001
        pass


def _set_hist(task_id, status, text=None):
    for h in _hist:
        if h.get("task_id") == task_id:
            h["status"] = status
            if text is not None:
                h["text"] = text
            break
    save_history()


def load_history():
    global _hist
    try:
        if os.path.exists(HISTORY_FILE):
            with open(HISTORY_FILE, encoding="utf-8") as f:
                _hist = json.load(f)
    except Exception:  # noqa: BLE001
        _hist = []
    # region agent log
    # 重启后仍在 running 的历史条目必然是孤儿（原进程已死），标记为 interrupted
    orphaned = [h.get("task_id") for h in _hist if h.get("status") == "running"]
    _dbg("load_history reconcile", "B", orphaned=orphaned)
    if orphaned:
        for h in _hist:
            if h.get("status") == "running":
                h["status"] = "interrupted"
        save_history()
    # endregion


# ---------------- 会话（多轮对话） ----------------

def _now_ts():
    return time.strftime("%Y-%m-%d %H:%M:%S")


def save_sessions():
    try:
        with open(SESSIONS_FILE, "w", encoding="utf-8") as f:
            json.dump(_sessions, f, ensure_ascii=False, indent=2)
    except Exception:  # noqa: BLE001
        pass


def load_sessions():
    global _sessions
    try:
        if os.path.exists(SESSIONS_FILE):
            with open(SESSIONS_FILE, encoding="utf-8") as f:
                _sessions = json.load(f)
    except Exception:  # noqa: BLE001
        _sessions = []


def _session_by_id(session_id):
    for s in _sessions:
        if s.get("session_id") == session_id:
            return s
    return None


def _new_session(session_id, agent, message):
    title = (message or "").strip().replace("\n", " ")[:40] or "新会话"
    s = {
        "session_id": session_id,
        "agent": agent,
        "agent_id": None,  # cursor 会话首轮结束后回填 agentId，用于续接
        "title": title,
        "created": _now_ts(),
        "updated": _now_ts(),
        "messages": [],
    }
    _sessions.insert(0, s)
    save_sessions()
    return s


def _append_session_msg(session_id, role, text, status="done"):
    s = _session_by_id(session_id)
    if not s:
        return
    s["messages"].append({
        "role": role,
        "text": text or "",
        "status": status,
        "ts": _now_ts(),
    })
    s["updated"] = _now_ts()
    save_sessions()


def _set_session_agent_id(session_id, agent_id):
    if not agent_id:
        return
    s = _session_by_id(session_id)
    if not s or s.get("agent_id"):
        return
    s["agent_id"] = agent_id
    save_sessions()


# ---------------- 文件系统工具 ----------------

def safe_resolve(rel_path):
    """把相对路径解析到允许的根目录内，返回绝对路径或 None。"""
    p = os.path.normpath(rel_path).lstrip(os.sep)
    if p in (".", ""):
        return os.path.abspath(TASKS_DIR) if os.path.isdir(TASKS_DIR) else None
    for root in ALLOWED_ROOTS:
        if not os.path.exists(root):
            continue
        cand = os.path.abspath(os.path.join(root, p))
        root_abs = os.path.abspath(root)
        if cand == root_abs or cand.startswith(root_abs + os.sep):
            return cand
    return None


def walk_tree(root, max_depth=3):
    if not os.path.isdir(root):
        return []
    items = []
    for entry in sorted(os.listdir(root), key=lambda x: (not os.path.isdir(os.path.join(root, x)), x.lower())):
        fp = os.path.join(root, entry)
        rel = os.path.relpath(fp, TASKS_DIR)
        # region agent log
        try:
            st = os.lstat(fp)
        except OSError as e:
            _dbg("walk_tree stat error", "A", path=fp, error=str(e))
            continue
        if stat.S_ISLNK(st.st_mode):
            _dbg("walk_tree skip symlink", "A", path=fp)
            continue
        # endregion
        if stat.S_ISDIR(st.st_mode):
            children = walk_tree(fp, max_depth - 1) if max_depth > 0 else []
            items.append({"name": entry, "type": "dir", "path": rel, "children": children})
        else:
            items.append({
                "name": entry, "type": "file", "path": rel,
                "size": st.st_size, "mtime": int(st.st_mtime),
                "kind": file_kind(entry),
            })
    # region agent log
    _dbg("walk_tree done", "A", root=root, count=len(items))
    # endregion
    return items


def latest_file(root):
    newest = None
    newest_mtime = -1
    if not os.path.isdir(root):
        return None
    for dirpath, _, filenames in os.walk(root):
        for name in filenames:
            fp = os.path.join(dirpath, name)
            try:
                st = os.stat(fp)
            except OSError:
                continue
            if st.st_mtime > newest_mtime:
                newest_mtime = st.st_mtime
                newest = {
                    "name": name,
                    "path": os.path.relpath(fp, TASKS_DIR),
                    "kind": file_kind(name),
                    "size": st.st_size,
                    "mtime": int(st.st_mtime),
                }
    return newest


def list_skills():
    if not os.path.isdir(SKILLS_DIR):
        return []
    out = []
    for entry in sorted(os.listdir(SKILLS_DIR)):
        skill_dir = os.path.join(SKILLS_DIR, entry)
        if not os.path.isdir(skill_dir):
            continue
        desc = ""
        smd = os.path.join(skill_dir, "SKILL.md")
        if os.path.exists(smd):
            try:
                with open(smd, encoding="utf-8") as f:
                    content = f.read(2000)
                for line in content.splitlines():
                    if line.lower().startswith("description:"):
                        desc = line.split(":", 1)[1].strip().strip('"').strip("'")[:120]
                        break
            except Exception:  # noqa: BLE001
                pass
        out.append({"name": entry, "description": desc})
    return out


def list_memory():
    if not os.path.isdir(MEMORY_DIR):
        return []
    out = []
    for entry in sorted(os.listdir(MEMORY_DIR), reverse=True):
        fp = os.path.join(MEMORY_DIR, entry)
        if os.path.isfile(fp):
            st = os.stat(fp)
            out.append({"name": entry, "size": st.st_size, "mtime": int(st.st_mtime)})
    return out


# ---------------- Office 预览 ----------------

NS_W = "{http://schemas.openxmlformats.org/wordprocessingml/2006/main}"
NS_A = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
NS_MAIN = "{http://schemas.openxmlformats.org/spreadsheetml/2006/main}"


def _xml_texts(root, tag):
    return [t.text or "" for t in root.iter(tag) if (t.text or "").strip()]


def preview_docx_xml(fp):
    with zipfile.ZipFile(fp) as z:
        xml = z.read("word/document.xml")
    root = ET.fromstring(xml)
    paras = []
    for p in root.iter(NS_W + "p"):
        line = "".join(t.text or "" for t in p.iter(NS_W + "t")).strip()
        if line:
            paras.append("<p>%s</p>" % htmlmod.escape(line))
    return "\n".join(paras) or "<p>（空文档）</p>"


def preview_docx(fp):
    if os.path.isfile(PANDOC):
        try:
            proc = subprocess.run(
                [PANDOC, fp, "-t", "html"],
                capture_output=True, timeout=40)
            if proc.returncode == 0 and proc.stdout:
                return proc.stdout.decode("utf-8", errors="replace")
        except Exception:  # noqa: BLE001
            pass
    return preview_docx_xml(fp)


def preview_pptx_xml(fp):
    slides = []
    with zipfile.ZipFile(fp) as z:
        names = sorted(
            n for n in z.namelist()
            if re.match(r"ppt/slides/slide\d+\.xml$", n)
        )
        for i, n in enumerate(names, 1):
            root = ET.fromstring(z.read(n))
            texts = [t for t in _xml_texts(root, NS_A + "t")]
            body = "<br>".join(htmlmod.escape(t) for t in texts) or "（空白页）"
            slides.append(
                "<section class='slide'><div class='slide-n'>SLIDE %02d</div>"
                "<div class='slide-body'>%s</div></section>" % (i, body)
            )
    return "".join(slides) or "<p>（空演示文稿）</p>"


def preview_pptx(fp):
    if os.path.isfile(VENV_PY):
        script = (
            "import sys\n"
            "from html import escape\n"
            "from pptx import Presentation\n"
            "prs = Presentation(sys.argv[1])\n"
            "out = []\n"
            "for i, slide in enumerate(prs.slides, 1):\n"
            "    bits = []\n"
            "    for shape in slide.shapes:\n"
            "        if hasattr(shape, 'text') and shape.text:\n"
            "            bits.append(escape(shape.text))\n"
            "    body = '<br>'.join(bits) or '（空白页）'\n"
            "    out.append('<section class=\"slide\"><div class=\"slide-n\">SLIDE %02d</div>"
            "<div class=\"slide-body\">%s</div></section>' % (i, body))\n"
            "print(''.join(out) or '<p>（空演示文稿）</p>')\n"
        )
        try:
            proc = subprocess.run(
                [VENV_PY, "-c", script, fp],
                capture_output=True, text=True, timeout=40)
            if proc.returncode == 0 and proc.stdout.strip():
                return proc.stdout
        except Exception:  # noqa: BLE001
            pass
    return preview_pptx_xml(fp)


def preview_xlsx(fp):
    if os.path.isfile(VENV_PY):
        script = (
            "import sys\n"
            "from html import escape\n"
            "import pandas as pd\n"
            "fp = sys.argv[1]\n"
            "xl = pd.ExcelFile(fp)\n"
            "parts = []\n"
            "for name in xl.sheet_names[:8]:\n"
            "    df = xl.parse(name, nrows=200)\n"
            "    parts.append('<h3>%s</h3>' % escape(str(name)))\n"
            "    parts.append(df.to_html(index=False, border=0, classes='sheet', max_cols=20))\n"
            "print(''.join(parts) or '<p>（空表格）</p>')\n"
        )
        try:
            proc = subprocess.run(
                [VENV_PY, "-c", script, fp],
                capture_output=True, text=True, timeout=40)
            if proc.returncode == 0 and proc.stdout.strip():
                return proc.stdout
        except Exception:  # noqa: BLE001
            pass
    return preview_xlsx_xml(fp)


def preview_xlsx_xml(fp):
    try:
        with zipfile.ZipFile(fp) as z:
            shared = []
            if "xl/sharedStrings.xml" in z.namelist():
                root = ET.fromstring(z.read("xl/sharedStrings.xml"))
                for si in root.iter(NS_MAIN + "si"):
                    shared.append("".join(t.text or "" for t in si.iter(NS_MAIN + "t")))
            sheet_name = next(
                (n for n in z.namelist() if n.startswith("xl/worksheets/sheet")),
                None,
            )
            if not sheet_name:
                return "<p>（无法解析表格）</p>"
            root = ET.fromstring(z.read(sheet_name))
            rows = []
            for row in list(root.iter(NS_MAIN + "row"))[:80]:
                cells = []
                for c in row.findall(NS_MAIN + "c"):
                    v = c.find(NS_MAIN + "v")
                    val = (v.text or "") if v is not None else ""
                    if c.get("t") == "s" and val.isdigit() and int(val) < len(shared):
                        val = shared[int(val)]
                    cells.append(htmlmod.escape(val))
                if any(cells):
                    rows.append("<tr>" + "".join("<td>%s</td>" % x for x in cells) + "</tr>")
            return "<table class='sheet'>" + "".join(rows) + "</table>" if rows else "<p>（空表格）</p>"
    except Exception as e:  # noqa: BLE001
        return "<p>无法预览：%s</p>" % htmlmod.escape(str(e))


def office_preview(fp):
    ext = os.path.splitext(fp)[1].lower()
    title = os.path.basename(fp)
    if ext in OFFICE_DOWNLOAD_EXTS:
        return {
            "kind": "office", "title": title, "download_only": True,
            "html": "<p class='hint'>旧版或 OpenDocument 文件请下载后用 Pages / WPS / Microsoft Office 打开。</p>",
        }
    if ext == ".docx":
        body = preview_docx(fp)
    elif ext == ".pptx":
        body = preview_pptx(fp)
    elif ext == ".xlsx":
        body = preview_xlsx(fp)
    else:
        return {"error": "unsupported office type", "kind": "office", "title": title}
    return {"kind": "office", "title": title, "html": body, "download_only": False}


def parse_range(header, size):
    if not header or not header.lower().startswith("bytes="):
        return None
    spec = header.split("=", 1)[1].split(",")[0].strip()
    if "-" not in spec:
        return None
    start_s, end_s = spec.split("-", 1)
    try:
        if start_s == "":
            n = int(end_s)
            start = max(0, size - n)
            end = size - 1
        else:
            start = int(start_s)
            end = int(end_s) if end_s else size - 1
    except ValueError:
        return None
    if start >= size or start < 0:
        return "unsatisfiable"
    end = min(end, size - 1)
    if end < start:
        return "unsatisfiable"
    return start, end


# ---------------- HTTP 处理 ----------------

class Handler(BaseHTTPRequestHandler):
    def _send(self, code, body, ctype="application/json; charset=utf-8"):
        if isinstance(body, str):
            body = body.encode("utf-8")
        self.send_response(code)
        self.send_header("Content-Type", ctype)
        self.send_header("Content-Length", str(len(body)))
        self.send_header("Cache-Control", "no-store")
        self.end_headers()
        self.wfile.write(body)

    def _json(self, code, obj):
        self._send(code, json.dumps(obj, ensure_ascii=False))

    def _disposition(self, name, inline=False):
        mode = "inline" if inline else "attachment"
        quoted = urllib.parse.quote(name, safe="")
        ascii_name = "file" + os.path.splitext(name)[1]
        ascii_name = "".join(c if ord(c) < 128 else "_" for c in ascii_name) or "file"
        return "%s; filename=\"%s\"; filename*=UTF-8''%s" % (mode, ascii_name, quoted)

    def _send_file_range(self, fp, ctype, download=False, head_only=False):
        size = os.path.getsize(fp)
        name = os.path.basename(fp)
        range_h = self.headers.get("Range") if not download else None
        parsed = parse_range(range_h, size) if range_h else None
        if parsed == "unsatisfiable":
            self.send_response(416)
            self.send_header("Content-Range", "bytes */%d" % size)
            self.send_header("Accept-Ranges", "bytes")
            self.end_headers()
            return
        if parsed:
            start, end = parsed
            length = end - start + 1
            self.send_response(206)
            self.send_header("Content-Type", ctype)
            self.send_header("Accept-Ranges", "bytes")
            self.send_header("Content-Range", "bytes %d-%d/%d" % (start, end, size))
            self.send_header("Content-Length", str(length))
            self.send_header("Cache-Control", "private, max-age=120")
            if download:
                self.send_header("Content-Disposition", self._disposition(name))
            self.end_headers()
            if head_only:
                return
            with open(fp, "rb") as f:
                f.seek(start)
                remaining = length
                while remaining > 0:
                    chunk = f.read(min(65536, remaining))
                    if not chunk:
                        break
                    self.wfile.write(chunk)
                    remaining -= len(chunk)
            return
        self.send_response(200)
        self.send_header("Content-Type", ctype)
        self.send_header("Accept-Ranges", "bytes")
        self.send_header("Content-Length", str(size))
        self.send_header("Cache-Control", "private, max-age=120")
        self.send_header("Content-Disposition", self._disposition(name, inline=not download))
        self.end_headers()
        if head_only:
            return
        with open(fp, "rb") as f:
            while True:
                chunk = f.read(65536)
                if not chunk:
                    break
                self.wfile.write(chunk)

    def do_OPTIONS(self):
        self.send_response(204)
        self.end_headers()

    def do_HEAD(self):
        self._handle_get(head_only=True)

    def do_GET(self):
        self._handle_get(head_only=False)

    def _handle_stream(self, tid):
        """SSE 流式输出：订阅任务事件队列，实时推送 status/tool/delta/done。"""
        q = queue.Queue()
        with _lock:
            _stream_queues.setdefault(tid, []).append(q)
            t = _tasks.get(tid)
            snapshot = dict(t) if t else {
                "task_id": tid, "status": "running", "text": "",
                "activity": "排队中…", "log": [], "agent": "cursor",
            }

        self.send_response(200)
        self.send_header("Content-Type", "text/event-stream; charset=utf-8")
        self.send_header("Cache-Control", "no-cache")
        self.send_header("Connection", "close")
        self.send_header("X-Accel-Buffering", "no")
        self.close_connection = True
        self.end_headers()

        def send_ev(obj):
            self.wfile.write(
                ("data: " + json.dumps(obj, ensure_ascii=False) + "\n\n").encode("utf-8"))
            self.wfile.flush()

        try:
            send_ev({
                "type": "snapshot",
                "text": snapshot.get("text") or "",
                "activity": snapshot.get("activity") or "",
                "status": snapshot.get("status"),
                "agent": snapshot.get("agent") or "cursor",
            })
            terminal_sent = False
            while not terminal_sent:
                try:
                    ev = q.get(timeout=0.5)
                except queue.Empty:
                    self.wfile.write(b": ping\n\n")
                    self.wfile.flush()
                    with _lock:
                        cur = _tasks.get(tid)
                    if cur and cur.get("status") in ("done", "error", "cancelled"):
                        st = cur["status"]
                        send_ev({
                            "type": st,
                            "text": cur.get("text") or "",
                            "error": cur.get("error"),
                        })
                        terminal_sent = True
                    continue
                send_ev(ev)
                if ev.get("type") in ("done", "error", "cancelled"):
                    terminal_sent = True
        except (BrokenPipeError, ConnectionResetError):
            pass
        finally:
            with _lock:
                subs = _stream_queues.get(tid)
                if subs and q in subs:
                    subs.remove(q)

    def _handle_get(self, head_only=False):
        parsed = urllib.parse.urlparse(self.path)
        path = parsed.path
        qs = urllib.parse.parse_qs(parsed.query)

        if path == "/" or path == "/index.html":
            return self._serve_static("index.html", "text/html; charset=utf-8")
        if path == "/marked.min.js":
            return self._serve_static("marked.min.js", "application/javascript; charset=utf-8")
        if path == "/favicon.svg":
            return self._serve_static("favicon.svg", "image/svg+xml")
        if path == "/favicon.ico":
            self.send_response(302)
            self.send_header("Location", "/favicon.svg")
            self.end_headers()
            return
        if path == "/manifest.webmanifest":
            return self._serve_static("manifest.webmanifest", "application/manifest+json")

        if path == "/api/health":
            ok = os.path.exists(os.path.join(HOME, ".openclaw", "openclaw.json"))
            return self._json(200, {
                "ok": ok, "tasks_dir": TASKS_DIR, "port": PORT,
                "agents": AGENTS, "version": VERSION, "product": PRODUCT,
            })

        if path == "/api/tasks":
            return self._json(200, walk_tree(TASKS_DIR))

        if path == "/api/latest":
            item = latest_file(TASKS_DIR)
            return self._json(200, item or {})

        if path == "/api/skills":
            return self._json(200, list_skills())

        if path == "/api/memory":
            return self._json(200, list_memory())

        if path == "/api/history":
            return self._json(200, _hist)

        if path == "/api/sessions":
            with _lock:
                out = [{
                    "session_id": s.get("session_id"),
                    "agent": s.get("agent"),
                    "agent_id": s.get("agent_id"),
                    "title": s.get("title"),
                    "created": s.get("created"),
                    "updated": s.get("updated"),
                    "message_count": len(s.get("messages") or []),
                } for s in _sessions]
            out.sort(key=lambda x: x.get("updated") or "", reverse=True)
            return self._json(200, out)

        if path.startswith("/api/session/"):
            sid = urllib.parse.unquote(path.split("/")[-1])
            with _lock:
                s = _session_by_id(sid)
                if not s:
                    return self._json(404, {"error": "session not found"})
                return self._json(200, dict(s))

        if path.startswith("/api/task/"):
            parts = path.split("/")
            # /api/task/<tid>/stream → SSE 流式输出
            if len(parts) >= 5 and parts[-1] == "stream":
                tid = urllib.parse.unquote(parts[-2])
                return self._handle_stream(tid)
            tid = urllib.parse.unquote(path.split("/")[-1])
            with _lock:
                t = _tasks.get(tid)
            if not t:
                return self._json(404, {"error": "task not found"})
            out = dict(t)
            out["task_id"] = tid
            return self._json(200, out)

        if path == "/api/file":
            rel = qs.get("path", [""])[0]
            fp = safe_resolve(rel)
            if not fp or not os.path.isfile(fp):
                return self._json(404, {"error": "file not found"})
            ext = os.path.splitext(fp)[1].lower()
            if ext not in TEXT_EXTS:
                return self._json(415, {"error": "binary file, use /api/media or /api/download"})
            with open(fp, encoding="utf-8", errors="replace") as f:
                return self._send(200, f.read()[:200000], "text/plain; charset=utf-8")

        if path == "/api/preview":
            rel = qs.get("path", [""])[0]
            fp = safe_resolve(rel)
            if not fp or not os.path.isfile(fp):
                return self._json(404, {"error": "file not found"})
            kind = file_kind(os.path.basename(fp))
            if kind != "office":
                return self._json(400, {"error": "not an office file", "kind": kind})
            try:
                return self._json(200, office_preview(fp))
            except Exception as e:  # noqa: BLE001
                return self._json(500, {"error": str(e), "kind": "office"})

        if path == "/api/media":
            rel = qs.get("path", [""])[0]
            fp = safe_resolve(rel)
            if not fp or not os.path.isfile(fp):
                return self._json(404, {"error": "file not found"})
            return self._send_file_range(
                fp, mime_for(os.path.basename(fp)), download=False, head_only=head_only)

        if path == "/api/download":
            rel = qs.get("path", [""])[0]
            fp = safe_resolve(rel)
            if not fp or not os.path.isfile(fp):
                return self._json(404, {"error": "file not found"})
            return self._send_file_range(
                fp, "application/octet-stream", download=True, head_only=head_only)

        return self._json(404, {"error": "not found"})

    def do_POST(self):
        parsed = urllib.parse.urlparse(self.path)
        path = parsed.path
        if path.startswith("/api/task/") and path.endswith("/cancel"):
            tid = urllib.parse.unquote(path.split("/")[-2])
            ok = cancel_task(tid)
            return self._json(200, {"ok": ok, "task_id": tid})
        if path.startswith("/api/session/") and path.endswith("/delete"):
            sid = urllib.parse.unquote(path.split("/")[-2])
            with _lock:
                before = len(_sessions)
                _sessions[:] = [s for s in _sessions if s.get("session_id") != sid]
                if len(_sessions) != before:
                    save_sessions()
            return self._json(200, {"ok": True})
        if path != "/api/task":
            return self._json(404, {"error": "not found"})
        length = int(self.headers.get("Content-Length", 0))
        if length > 200000:
            return self._json(413, {"error": "message too long"})
        try:
            data = json.loads(self.rfile.read(length) or b"{}")
        except Exception:  # noqa: BLE001
            return self._json(400, {"error": "bad json"})
        message = (data.get("message") or "").strip()
        agent = data.get("agent") or "main"
        if not message:
            return self._json(400, {"error": "empty message"})
        if agent not in AGENTS:
            return self._json(400, {"error": "unknown agent"})

        # 会话：支持继续对话。前端传 sessionId 则复用，否则新建会话。
        session_id = (data.get("sessionId") or "").strip()
        with _lock:
            if not (session_id and _session_by_id(session_id)):
                session_id = "s-%d-%d" % (int(time.time() * 1000), len(_sessions) + 1)
                _new_session(session_id, agent, message)
            _append_session_msg(session_id, "user", message, "done")

            task_id = "%s-%d" % (int(time.time() * 1000), len(_tasks) + 1)
            _tasks[task_id] = {
                "task_id": task_id,
                "status": "running", "text": "", "error": None,
                "agent": agent, "message": message,
                "session_id": session_id,
                "started": time.time(), "duration": 0,
                "activity": "排队中…", "log": [],
            }
            _hist.insert(0, {"task_id": task_id, "agent": agent, "message": message,
                             "status": "running", "session_id": session_id,
                             "time": time.strftime("%Y-%m-%d %H:%M:%S")})
            save_history()
        target = run_cursor_agent if agent == "cursor" else run_openclaw_agent
        threading.Thread(target=target, args=(task_id, agent, message, session_id), daemon=True).start()
        return self._json(200, {"task_id": task_id, "session_id": session_id})

    def _serve_static(self, name, ctype):
        fp = os.path.join(WORKBENCH_DIR, name)
        if not os.path.exists(fp):
            return self._json(404, {"error": "missing %s" % name})
        with open(fp, encoding="utf-8") as f:
            return self._send(200, f.read(), ctype)

    def log_message(self, *args):
        pass


def main():
    load_history()
    load_sessions()
    os.makedirs(TASKS_DIR, exist_ok=True)
    log("starting %s v%s on http://%s:%d" % (PRODUCT, VERSION, HOST, PORT))
    log("tasks dir: %s" % TASKS_DIR)
    srv = ThreadingHTTPServer((HOST, PORT), Handler)
    srv.serve_forever()


if __name__ == "__main__":
    main()
